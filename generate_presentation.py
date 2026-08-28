import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    base_dir = r'c:\Users\Josué\Documents\App Rafael'

    # Imagens
    img_mobile_compras = os.path.join(base_dir, 'app_compras_mobile_ui_1787233856843.jpg')
    img_mobile_separacao = os.path.join(base_dir, 'app_separacao_mobile_ui_1787233871196.jpg')
    img_web_desktop = os.path.join(base_dir, 'sistema_web_desktop_ui_1787233888011.jpg')

    # Paleta de Cores Padrão App Érica
    C_PRIMARY = RGBColor(16, 185, 129)       # Emerald Green #10B981
    C_PRIMARY_DARK = RGBColor(5, 150, 105)   # #059669
    C_DARK_BG = RGBColor(15, 23, 42)         # #0F172A
    C_DARK_SURFACE = RGBColor(30, 41, 59)    # #1E293B
    C_LIGHT_BG = RGBColor(248, 250, 252)     # #F8FAFC
    C_WHITE = RGBColor(255, 255, 255)
    C_TEXT_MAIN = RGBColor(15, 23, 42)
    C_TEXT_MUTED = RGBColor(100, 116, 139)
    C_GOLD = RGBColor(245, 158, 11)          # #F59E0B
    C_CARD_BORDER = RGBColor(226, 232, 240)
    C_LIGHT_GREEN = RGBColor(240, 253, 244)  # #F0FDF4
    C_RED = RGBColor(220, 38, 38)

    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title, category, subtitle=None, dark=False):
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_tag = tb_tag.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
        pt = tf_tag.paragraphs[0]
        pt.text = category.upper()
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = C_PRIMARY if dark else C_PRIMARY_DARK
        
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.55))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        ptt = tf_title.paragraphs[0]
        ptt.text = title
        ptt.font.size = Pt(22)
        ptt.font.bold = True
        ptt.font.color.rgb = C_WHITE if dark else C_TEXT_MAIN

        if subtitle:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.4))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            ps = tf_sub.paragraphs[0]
            ps.text = subtitle
            ps.font.size = Pt(12)
            ps.font.color.rgb = RGBColor(148, 163, 184) if dark else C_TEXT_MUTED

    # =============================================================
    # SLIDE 1: CAPA (DARK THEME)
    # =============================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, C_DARK_BG)

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(4.8), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_DARK_SURFACE
    badge.line.color.rgb = C_PRIMARY
    tf_b = badge.text_frame
    tf_b.margin_top = Inches(0.06)
    pb = tf_b.paragraphs[0]
    pb.text = '✦  ECOSSISTEMA ALS 10 & CONECTA'
    pb.font.size = Pt(11)
    pb.font.bold = True
    pb.font.color.rgb = C_PRIMARY

    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = 'APP & SISTEMA DE COMPRAS E RATEIO\nDA PLANILHA MATRIZ'
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf1.add_paragraph()
    p2.text = 'Sua planilha atual transformada em um aplicativo simples, ágil e com visual idêntico ao App Érica, resolvendo os problemas do Excel.'
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(10)

    cards_s1 = [
        ('1. Pedido de Compra', 'Mesmos campos da aba "ALS 10 BAZAR", com cálculo de caixas e PDF direto no WhatsApp.'),
        ('2. Limite de Preço', 'Mesmas fórmulas de ICMS, Créditos e Despesas do PDV, calculadas na hora sem risco de erro.'),
        ('3. Rateio das 20 Lojas', 'Mesmos grupos A, B e C, mas arredondando peças automaticamente sem decimais quebrados.'),
        ('4. Simples & Econômico', 'Sem servidores caros ou mensalidades pesadas: focado exatamente no que você precisa hoje.')
    ]

    for idx, (t, d) in enumerate(cards_s1):
        cx = Inches(0.8 + idx * 2.95)
        cshape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(4.7), Inches(2.8), Inches(2.1))
        cshape.fill.solid()
        cshape.fill.fore_color.rgb = C_DARK_SURFACE
        cshape.line.color.rgb = RGBColor(51, 65, 85)
        
        tfc = cshape.text_frame
        tfc.word_wrap = True
        tfc.margin_left = tfc.margin_right = Inches(0.15)
        tfc.margin_top = Inches(0.15)
        
        pc1 = tfc.paragraphs[0]
        pc1.text = t
        pc1.font.size = Pt(13)
        pc1.font.bold = True
        pc1.font.color.rgb = C_PRIMARY
        
        pc2 = tfc.add_paragraph()
        pc2.text = d
        pc2.font.size = Pt(10.5)
        pc2.font.color.rgb = RGBColor(203, 213, 225)
        pc2.space_before = Pt(6)

    # =============================================================
    # SLIDE 2: O QUE RESOLVEMOS DA SUA PLANILHA
    # =============================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, C_LIGHT_BG)
    add_header(s2, 'Mantendo Tudo o que Funciona e Eliminando as Dores', 'FOCO NA PRATICIDADE', 'Mantemos 100% da sua regra de negócio atual, corrigindo apenas o que atrasa o dia a dia.')

    pains = [
        ('❌ Fim dos Números Decimais no Rateio', 'Na planilha atual, a fórmula gera frações (ex: 7,3 peças por loja). No novo app, o cálculo distribui o lote exato em números inteiros sem sobras.'),
        ('❌ Fim das Fórmulas que Apagam ou Quebram', 'Inserir uma nova linha de produto não vai mais desconfigurar os totais. As fórmulas ficam protegidas e automáticas no app.'),
        ('❌ Fim da Redigitação Repetitiva', 'Os nomes dos produtos, referências e fornecedores ficam salvos na memória. Ao começar a digitar, o sistema já preenche sozinho.'),
        ('❌ Envio em PDF com 1 Clique', 'Chega de desalinhar na impressão ou tirar print da tela. O app gera o pedido formatado e pronto para mandar ao fornecedor no WhatsApp.')
    ]

    for idx, (title_p, desc_p) in enumerate(pains):
        row = idx // 2
        col = idx % 2
        x = Inches(0.8 + col * 5.95)
        y = Inches(1.9 + row * 2.4)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(2.15))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_CARD_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title_p
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc_p
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(8)

    # =============================================================
    # SLIDE 3: VISUAL DO APP MOBILE (COM IMAGENS)
    # =============================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, C_LIGHT_BG)
    add_header(s3, 'Visual do Aplicativo Mobile (Padrão App Érica)', 'EXPERIÊNCIA MOBILE', 'Design moderno com verde esmeralda, cards intuitivos e operação ultrarrápida no celular.')

    # Left Image (Mobile Compras)
    if os.path.exists(img_mobile_compras):
        s3.shapes.add_picture(img_mobile_compras, Inches(0.8), Inches(1.85), width=Inches(2.8))

    # Center Image (Mobile Separação)
    if os.path.exists(img_mobile_separacao):
        s3.shapes.add_picture(img_mobile_separacao, Inches(3.85), Inches(1.85), width=Inches(2.8))

    # Right Content Card
    rc_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.85), Inches(5.6), Inches(4.9))
    rc_card.fill.solid()
    rc_card.fill.fore_color.rgb = C_WHITE
    rc_card.line.color.rgb = C_PRIMARY
    tf_rc = rc_card.text_frame
    tf_rc.word_wrap = True
    tf_rc.margin_left = tf_rc.margin_right = Inches(0.25)
    tf_rc.margin_top = Inches(0.2)

    prc1 = tf_rc.paragraphs[0]
    prc1.text = '✦ Funcionalidades na Tela do Celular:'
    prc1.font.size = Pt(15)
    prc1.font.bold = True
    prc1.font.color.rgb = C_TEXT_MAIN

    bullets_app = [
        'Pedido de Compra Ágil: Lançamento de fornecedor, caixas e unidades em segundos.',
        'Simulador Fiscal em Tempo Real: Mostra na hora o custo real e se a margem PDV está aprovada (em verde).',
        'Romaneio de Separação das 20 Lojas: Grade visual dividida nos Grupos A, B e C com quantidades exatas sem decimais.',
        'Zero Travamentos: Funciona rápido mesmo com sinal oscilante no depósito.'
    ]
    for b in bullets_app:
        pb = tf_rc.add_paragraph()
        pb.text = f'✓ {b}'
        pb.font.size = Pt(11.5)
        pb.font.color.rgb = C_TEXT_MUTED
        pb.space_before = Pt(8)

    # =============================================================
    # SLIDE 4: VISUAL DO SISTEMA WEB DESKTOP (COM IMAGEM)
    # =============================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, C_LIGHT_BG)
    add_header(s4, 'Visual do Sistema no Computador (Web / Desktop)', 'VISÃO NO ESCRITÓRIO', 'Tela ampla e limpa para emissão rápida de compras e conferência de limites fiscais.')

    # Large Dashboard Image
    if os.path.exists(img_web_desktop):
        s4.shapes.add_picture(img_web_desktop, Inches(0.8), Inches(1.85), width=Inches(7.8))

    # Right Card with Features
    w_card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.85), Inches(3.7), Inches(4.4))
    w_card.fill.solid()
    w_card.fill.fore_color.rgb = C_WHITE
    w_card.line.color.rgb = C_CARD_BORDER
    tf_w = w_card.text_frame
    tf_w.word_wrap = True
    tf_w.margin_left = tf_w.margin_right = Inches(0.2)
    tf_w.margin_top = Inches(0.2)

    pw = tf_w.paragraphs[0]
    pw.text = 'Praticidade no PC:'
    pw.font.size = Pt(14)
    pw.font.bold = True
    pw.font.color.rgb = C_PRIMARY_DARK

    w_items = [
        'Tabela de itens idêntica à planilha, com preenchimento em lote.',
        'Auditoria fiscal integrada na barra lateral (ICMS 11%, Crédito 19.5%).',
        'Cálculo de margem e preço teto atualizado a cada item digitado.',
        'Exportação de espelho de pedido em PDF oficial para o fornecedor.'
    ]
    for wi in w_items:
        pwi = tf_w.add_paragraph()
        pwi.text = f'• {wi}'
        pwi.font.size = Pt(11)
        pwi.font.color.rgb = C_TEXT_MUTED
        pwi.space_before = Pt(8)

    # =============================================================
    # SLIDE 5: COMPARAÇÃO DIRETA COM A PLANILHA
    # =============================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, C_LIGHT_BG)
    add_header(s5, 'Quadro Comparativo: Planilha Matriz vs. Novo App', 'BENEFÍCIOS REAIS', 'A mesma lógica que você já domina, com muito mais segurança e rapidez.')

    rows = [
        ('Etapa do Processo', 'Como é na Planilha Hoje', 'Como fica no Novo App'),
        ('Emissão do Pedido', 'Digita tudo do zero toda vez no Excel.', 'Preenche rápido e memoriza produtos e fornecedores.'),
        ('Multiplicador de Caixas', 'Risco de errar ou apagar a fórmula F×G.', 'Cálculo travado e 100% automático.'),
        ('Limite de Preço / Impostos', 'Aba separada com links que podem sumir.', 'Auditoria instantânea de margem no mesmo painel.'),
        ('Rateio das 20 Lojas', 'Dá frações (7,3 peças) e exige ajuste manual.', 'Rateia na hora em números inteiros sem sobras.'),
        ('Envio ao Fornecedor', 'Desformata na impressão ou envio.', 'Gera PDF oficial pronto para WhatsApp com 1 clique.')
    ]

    table_shape = s5.shapes.add_table(6, 3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.9))
    table = table_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(4.45)
    table.columns[2].width = Inches(4.45)

    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11 if r_idx > 0 else 12)
            p.font.bold = (r_idx == 0 or c_idx == 0)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_DARK_BG
                p.font.color.rgb = C_WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if r_idx % 2 == 1 else C_LIGHT_GREEN
                if c_idx == 2:
                    p.font.color.rgb = C_PRIMARY_DARK
                    p.font.bold = True
                elif c_idx == 1:
                    p.font.color.rgb = RGBColor(185, 28, 28)
                else:
                    p.font.color.rgb = C_TEXT_MAIN

    # =============================================================
    # SLIDE 6: CRONOGRAMA RÁPIDO (2 A 3 SEMANAS)
    # =============================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, C_LIGHT_BG)
    add_header(s6, 'Entrega Rápida e Implantação Sem Burocracia', 'CRONOGRAMA ENXUTO', 'Desenvolvimento ágil com entrega total entre 15 e 20 dias corridos.')

    steps_lean = [
        ('Semana 1', 'Construção da Base & Visual', 'Configuração das regras da planilha matriz, parametrização dos impostos e aplicação do design App Érica.'),
        ('Semana 2', 'Pedidos & Rateio 20 Lojas', 'Módulo de lançamento de pedidos, cálculo de caixas, limite de preço e motor de rateio sem decimais.'),
        ('Semana 3', 'Testes, PDF & Entrega', 'Geração de PDF para WhatsApp, validação com 1 pedido real junto com o Rafael e ativação.')
    ]

    for idx, (sem, tit, desc) in enumerate(steps_lean):
        x = Inches(0.8 + idx * 3.95)
        y = Inches(2.1)
        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(4.5))
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        c.line.color.rgb = C_PRIMARY_DARK if idx == 2 else C_CARD_BORDER
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.25)
        
        p_num = tf.paragraphs[0]
        p_num.text = sem.upper()
        p_num.font.size = Pt(12)
        p_num.font.bold = True
        p_num.font.color.rgb = C_PRIMARY_DARK
        
        p_t = tf.add_paragraph()
        p_t.text = tit
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_MAIN
        p_t.space_before = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(10)

    # =============================================================
    # SLIDE 7: PROPOSTA COMERCIAL ECONÔMICA (DARK THEME)
    # =============================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, C_DARK_BG)
    add_header(s7, 'Investimento Acessível & Direto ao Ponto', 'PROPOSTA ECONÔMICA', 'Projeto sob medida com excelente custo-benefício para substituir a planilha de vez.', dark=True)

    props = [
        ('Opção 1: À Vista (com Desconto Especial)', 'R$ 1.900,00', [
            '50% de entrada no início do projeto (R$ 950)',
            '50% na entrega e aprovação final (R$ 950)',
            'Maior economia e prioridade total na fila de entrega',
            'Garantia e suporte assistido na primeira emissão de pedido'
        ], C_PRIMARY),
        ('Opção 2: Parcelado em 3x Sem Juros', '3x de R$ 750,00', [
            'Entrada de R$ 750,00 no início do projeto',
            '2ª Parcela de R$ 750,00 em 30 dias',
            '3ª Parcela de R$ 750,00 em 60 dias',
            'Total parcelado: R$ 2.250,00 com investimento diluído'
        ], C_GOLD)
    ]

    for idx, (pname, pval, pitems, pcolor) in enumerate(props):
        x = Inches(0.8 + idx * 5.95)
        y = Inches(1.9)
        c = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(4.9))
        c.fill.solid()
        c.fill.fore_color.rgb = C_DARK_SURFACE
        c.line.color.rgb = pcolor
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = pname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        pv = tf.add_paragraph()
        pv.text = pval
        pv.font.size = Pt(28)
        pv.font.bold = True
        pv.font.color.rgb = pcolor
        pv.space_before = Pt(4)
        pv.space_after = Pt(12)
        
        for item in pitems:
            pi = tf.add_paragraph()
            pi.text = f'✦  {item}'
            pi.font.size = Pt(12)
            pi.font.color.rgb = RGBColor(203, 213, 225)
            pi.space_before = Pt(6)

    output_final = r'c:\Users\Josué\Documents\App Rafael\Apresentacao_App_Compras_Matriz_Final.pptx'
    prs.save(output_final)
    print(f'Sucesso: Apresentação salva em {output_final}')

if __name__ == '__main__':
    create_presentation()

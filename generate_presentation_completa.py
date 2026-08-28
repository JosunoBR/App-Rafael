import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_complete_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    base_dir = r'c:\Users\Josué\Documents\App Rafael'

    # Imagens dos mockups
    img_mobile_compras = os.path.join(base_dir, 'app_compras_mobile_ui_1787233856843.jpg')
    img_mobile_separacao = os.path.join(base_dir, 'app_separacao_mobile_ui_1787233871196.jpg')
    img_web_desktop = os.path.join(base_dir, 'sistema_web_desktop_ui_1787233888011.jpg')

    # Paleta de Cores Padrão App Érica
    C_PRIMARY = RGBColor(16, 185, 129)       # Emerald Green #10B981
    C_PRIMARY_DARK = RGBColor(5, 150, 105)   # #059669
    C_DARK_BG = RGBColor(15, 23, 42)         # #0F172A
    C_DARK_SURFACE = RGBColor(30, 41, 59)    # #1E293B
    C_LIGHT_BG = RGBColor(248, 250, 252)     # #F8FAFC
    C_WHITE = RGBColor(255, 255, 255)
    C_TEXT_MAIN = RGBColor(15, 23, 42)
    C_TEXT_MUTED = RGBColor(100, 116, 139)
    C_GOLD = RGBColor(245, 158, 11)          # #F59E0B
    C_CARD_BORDER = RGBColor(226, 232, 240)
    C_LIGHT_GREEN = RGBColor(240, 253, 244)  # #F0FDF4
    C_RED = RGBColor(220, 38, 38)

    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title, category, subtitle=None, dark=False):
        tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_tag = tb_tag.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
        pt = tf_tag.paragraphs[0]
        pt.text = category.upper()
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = C_PRIMARY if dark else C_PRIMARY_DARK
        
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.55))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        ptt = tf_title.paragraphs[0]
        ptt.text = title
        ptt.font.size = Pt(22)
        ptt.font.bold = True
        ptt.font.color.rgb = C_WHITE if dark else C_TEXT_MAIN

        if subtitle:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.4))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            ps = tf_sub.paragraphs[0]
            ps.text = subtitle
            ps.font.size = Pt(12)
            ps.font.color.rgb = RGBColor(148, 163, 184) if dark else C_TEXT_MUTED

    # =============================================================
    # SLIDE 1: CAPA OFICIAL (DARK THEME)
    # =============================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, C_DARK_BG)

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(4.8), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_DARK_SURFACE
    badge.line.color.rgb = C_PRIMARY
    tf_b = badge.text_frame
    tf_b.margin_top = Inches(0.06)
    pb = tf_b.paragraphs[0]
    pb.text = '✦  ECOSSISTEMA INTEGRADO ALS 10 & CONECTA'
    pb.font.size = Pt(11)
    pb.font.bold = True
    pb.font.color.rgb = C_PRIMARY

    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = 'SISTEMA INTEGRADO DE COMPRAS,\nLIMITE DE PREÇO & SEPARAÇÃO'
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf1.add_paragraph()
    p2.text = 'Plataforma Corporativa Completa: Portal Web de Gestão + App Mobile Multifunção com a identidade visual do App Érica.'
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(10)

    cards_s1 = [
        ('1. Gestão Central Web', 'Portal para PC com emissão de pedidos, cadastro de itens e controle financeiro.'),
        ('2. Engenharia Fiscal', 'Auditoria tributária em tempo real (ICMS, Créditos, Despesas e Margem PDV).'),
        ('3. Rateio das 20 Lojas', 'Distribuição matemática ponderada sem decimais quebrados para os grupos A, B e C.'),
        ('4. App Mobile Completo', 'Mobilidade total para o comprador em viagens e romaneio digital sem papel no CD.')
    ]

    for idx, (t, d) in enumerate(cards_s1):
        cx = Inches(0.8 + idx * 2.95)
        cshape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(4.7), Inches(2.8), Inches(2.1))
        cshape.fill.solid()
        cshape.fill.fore_color.rgb = C_DARK_SURFACE
        cshape.line.color.rgb = RGBColor(51, 65, 85)
        
        tfc = cshape.text_frame
        tfc.word_wrap = True
        tfc.margin_left = tfc.margin_right = Inches(0.15)
        tfc.margin_top = Inches(0.15)
        
        pc1 = tfc.paragraphs[0]
        pc1.text = t
        pc1.font.size = Pt(13)
        pc1.font.bold = True
        pc1.font.color.rgb = C_PRIMARY
        
        pc2 = tfc.add_paragraph()
        pc2.text = d
        pc2.font.size = Pt(10.5)
        pc2.font.color.rgb = RGBColor(203, 213, 225)
        pc2.space_before = Pt(6)

    # =============================================================
    # SLIDE 2: O CENÁRIO ATUAL (DIAGNÓSTICO DA PLANILHA MATRIZ)
    # =============================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, C_LIGHT_BG)
    add_header(s2, 'Diagnóstico da Operação Atual (Planilha MATRIZ.xlsx)', 'O DESAFIO ATUAL', 'Identificação dos gargalos operacionais no modelo de gestão baseado em planilhas.')

    pains = [
        ('⚠️ Fórmulas e Links Quebradiços', 'A planilha possui interligações manuais entre abas. Inserir novas linhas de produtos frequentemente corrompe fórmulas de totalização e rateio.'),
        ('🔢 Rateio com Decimais & Ajuste Manual', 'A fórmula matemática divide peças em números quebrados (ex: 7,3 peças por loja), forçando o comprador a recalcular e arredondar loja a loja na mão.'),
        ('📋 Retrabalho de Digitação Contínuo', 'Sem cadastro unificado de itens ou fornecedores, é necessário redigitar referências, descrições e dados fiscais a cada novo pedido emitido.'),
        ('📄 Papel & Falta de Mobilidade no CD', 'A separação das mercadorias depende de impressão em papel ou de levar notebook para o galpão, gerando riscos de erros no despacho.')
    ]

    for idx, (title_p, desc_p) in enumerate(pains):
        row = idx // 2
        col = idx % 2
        x = Inches(0.8 + col * 5.95)
        y = Inches(1.9 + row * 2.4)
        
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(2.15))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_CARD_BORDER
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title_p
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_RED
        
        p2 = tf.add_paragraph()
        p2.text = desc_p
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(8)

    # =============================================================
    # SLIDE 3: A SOLUÇÃO INTEGRADA (PORTAL WEB + APP)
    # =============================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, C_LIGHT_BG)
    add_header(s3, 'A Nova Plataforma Integrada: Gestão no PC e Celular', 'ARQUITETURA DA SOLUÇÃO', 'Ambiente único em nuvem conectando o escritório de compras, viagens e a expedição do CD.')

    sol_cards = [
        ('🖥️ Portal Web do Comprador & Gestor', 'Módulo Desktop / Navegador (PC do Escritório)', [
            'Cadastro mestre de Fornecedores e Catálogo de Produtos',
            'Geração de Pedidos com embalagens fechadas e custos',
            'Simulador fiscal automático de Custo Real e Limite de Preço',
            'Motor de Rateio Inteligente para 20 lojas sem sobras',
            'Exportação de Pedidos em PDF padronizado com 1 clique'
        ]),
        ('📱 App Mobile Multifunção (Padrão App Érica)', 'Módulo Celular / Tablet Android (Gestão & Galpão)', [
            'Lançamento de compras direto do celular em feiras/viagens',
            'Consulta de histórico de compras na mesa do fornecedor',
            'Romaneio digital por filial na tela do operador no CD',
            'Checklist de separação por volume/caixa sem papel',
            'Funciona 100% sem internet no estoque (zero travamentos)'
        ])
    ]

    for idx, (stitle, ssub, bullets) in enumerate(sol_cards):
        x = Inches(0.8 + idx * 5.95)
        y = Inches(1.9)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_PRIMARY_DARK if idx == 0 else C_PRIMARY
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_MAIN
        
        p_sub = tf.add_paragraph()
        p_sub.text = ssub
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_PRIMARY_DARK
        p_sub.space_before = Pt(2)
        p_sub.space_after = Pt(12)
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f'✓  {b}'
            pb.font.size = Pt(11.5)
            pb.font.color.rgb = C_TEXT_MAIN
            pb.space_before = Pt(6)

    # =============================================================
    # SLIDE 4: MÓDULO WEB DESKTOP (COM MOCKUP VISUAL)
    # =============================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, C_LIGHT_BG)
    add_header(s4, 'Módulo Web: Painel de Pedidos & Engenharia Fiscal', 'PORTAL DO COMPRADOR', 'Interface desktop com tabela ampla para digitação rápida e auditoria fiscal integrada.')

    if os.path.exists(img_web_desktop):
        s4.shapes.add_picture(img_web_desktop, Inches(0.8), Inches(1.85), width=Inches(7.8))

    w_card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.85), Inches(3.7), Inches(4.4))
    w_card.fill.solid()
    w_card.fill.fore_color.rgb = C_WHITE
    w_card.line.color.rgb = C_PRIMARY
    tf_w = w_card.text_frame
    tf_w.word_wrap = True
    tf_w.margin_left = tf_w.margin_right = Inches(0.2)
    tf_w.margin_top = Inches(0.2)

    pw = tf_w.paragraphs[0]
    pw.text = 'Diferenciais do Portal:'
    pw.font.size = Pt(14)
    pw.font.bold = True
    pw.font.color.rgb = C_PRIMARY_DARK

    w_items = [
        'Multiplicador de Caixas: Qtd no Pec × Qtd Pec = Qtd Total calculado na hora.',
        'Auditoria Fiscal Lateral: ICMS (11%), Crédito (19.5%) e Custo Efetivo instantâneo.',
        'Selo de Margem: Alerta visual verde confirmando viabilidade de venda no PDV.',
        'Exportação Direta: Espelho do pedido em PDF padronizado para o fornecedor.'
    ]
    for wi in w_items:
        pwi = tf_w.add_paragraph()
        pwi.text = f'✓ {wi}'
        pwi.font.size = Pt(11)
        pwi.font.color.rgb = C_TEXT_MUTED
        pwi.space_before = Pt(8)

    # =============================================================
    # SLIDE 5: LIMITE DE PREÇO & FORMAÇÃO DE CUSTO REAL
    # =============================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, C_LIGHT_BG)
    add_header(s5, 'Engenharia Tributária: Custo Real & Margem Líquida', 'LIMITE DE PREÇO', 'Parametrização automática das regras fiscais para blindar a lucratividade do negócio.')

    eq_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.45))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = C_LIGHT_GREEN
    eq_box.line.color.rgb = C_PRIMARY
    tfe = eq_box.text_frame
    tfe.word_wrap = True
    tfe.margin_left = tfe.margin_right = Inches(0.3)
    tfe.margin_top = Inches(0.15)
    pe1 = tfe.paragraphs[0]
    pe1.text = 'FÓRMULA AUTOMÁTICA DE CUSTO REAL EFETIVO:'
    pe1.font.size = Pt(11)
    pe1.font.bold = True
    pe1.font.color.rgb = C_PRIMARY_DARK
    pe2 = tfe.add_paragraph()
    pe2.text = 'Custo Total Produto = Valor Compra + (PDV × % Despesas Fixas/Impostos) - (Valor Compra × % Crédito ICMS)'
    pe2.font.size = Pt(13.5)
    pe2.font.bold = True
    pe2.font.color.rgb = C_TEXT_MAIN
    pe2.space_before = Pt(4)

    sub_items = [
        ('🛡️ Parametrização Tributária', 'Configuração de ICMS (11%), PIS/COFINS (3%), Custos Fixos (26%) e Crédito de Entrada (19,5%) sem risco de fórmulas desconfiguradas.'),
        ('🎯 Validador de Margem PDV', 'Alerta visual instantâneo: Verde (Margem Saudável), Amarelo (Margem Apertada) e Vermelho (Preço de compra inviabiliza a venda no teto do PDV).'),
        ('📊 Apoio à Negociação', 'O comprador sabe exatamente até quanto pode pagar no fornecedor para manter a lucratividade padrão da rede.')
    ]

    for idx, (t, d) in enumerate(sub_items):
        x = Inches(0.8 + idx * 3.95)
        y = Inches(3.6)
        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(3.2))
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        c.line.color.rgb = C_CARD_BORDER
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY_DARK
        
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(8)

    # =============================================================
    # SLIDE 6: MATRIZ DE RATEIO PONDERADO (20 LOJAS)
    # =============================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, C_LIGHT_BG)
    add_header(s6, 'Motor de Rateio Inteligente para 20 Lojas (Sem Decimais)', 'RATEIO AUTOMÁTICO', 'Substitui a "Calculadora", distribuindo produtos pelos pesos ponderados sem frações.')

    clusters = [
        ('GRUPO A (8 LOJAS)', 'Peso: 2.5 | 20 Pontos Ponderados', [
            'Ponta Grossa Centro', 'Reserva', 'Tibagi', 'Nova Rússia',
            'Javert', 'Ivaí', 'Irati Centro', 'Campo Largo'
        ], RGBColor(220, 38, 38)),
        ('GRUPO B (8 LOJAS)', 'Peso: 1.75 | 14 Pontos Ponderados', [
            'Castro', 'Imbituva', 'Santa Paula', 'Prudentópolis',
            'Guarapuava', 'Imbaú', 'Rio Azul', 'Rebouças'
        ], C_GOLD),
        ('GRUPO C (4 LOJAS)', 'Peso: 1.25 | 5 Pontos Ponderados', [
            'Depósito Central', 'Teixeira Soares', 'Mallet', 'Ipiranga'
        ], C_PRIMARY_DARK)
    ]

    for idx, (cname, cdetail, stores, color) in enumerate(clusters):
        x = Inches(0.8 + idx * 3.95)
        y = Inches(1.9)
        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(4.9))
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        c.line.color.rgb = color
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = cname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_sub = tf.add_paragraph()
        p_sub.text = cdetail
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_TEXT_MUTED
        p_sub.space_before = Pt(2)
        p_sub.space_after = Pt(10)
        
        for st in stores:
            pst = tf.add_paragraph()
            pst.text = f'• {st}'
            pst.font.size = Pt(11)
            pst.font.color.rgb = C_TEXT_MAIN
            pst.space_before = Pt(4)

    # =============================================================
    # SLIDE 7: APP MOBILE (COM MOCKUPS VISUAIS)
    # =============================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, C_LIGHT_BG)
    add_header(s7, 'App Mobile Multifunção: Compras em Viagens & Romaneio no CD', 'APLICATIVO MOBILE', 'Visual idêntico ao App Érica: Emerald Green, botões rápidos, intuitivo e com modo offline.')

    # Left Image (Mobile Compras)
    if os.path.exists(img_mobile_compras):
        s7.shapes.add_picture(img_mobile_compras, Inches(0.8), Inches(1.85), width=Inches(2.8))

    # Center Image (Mobile Separação)
    if os.path.exists(img_mobile_separacao):
        s7.shapes.add_picture(img_mobile_separacao, Inches(3.85), Inches(1.85), width=Inches(2.8))

    # Right Content Card
    rc_card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.85), Inches(5.6), Inches(4.9))
    rc_card.fill.solid()
    rc_card.fill.fore_color.rgb = C_WHITE
    rc_card.line.color.rgb = C_PRIMARY
    tf_rc = rc_card.text_frame
    tf_rc.word_wrap = True
    tf_rc.margin_left = tf_rc.margin_right = Inches(0.25)
    tf_rc.margin_top = Inches(0.2)

    prc1 = tf_rc.paragraphs[0]
    prc1.text = '✦ Funcionalidades no Celular / Tablet:'
    prc1.font.size = Pt(15)
    prc1.font.bold = True
    prc1.font.color.rgb = C_TEXT_MAIN

    bullets_app = [
        'Compras em Viagens: O Rafael lança pedidos e simula margens direto na mesa do fornecedor em feiras.',
        'Envio por WhatsApp: Compartilha o PDF do pedido no mesmo instante com o representante.',
        'Romaneio Digital no CD: O separador visualiza a quantidade de cada loja sem usar prancheta de papel.',
        'Zero Travamentos: Funciona 100% offline no depósito e sincroniza sozinho quando tiver internet.'
    ]
    for b in bullets_app:
        pb = tf_rc.add_paragraph()
        pb.text = f'✓ {b}'
        pb.font.size = Pt(11.5)
        pb.font.color.rgb = C_TEXT_MUTED
        pb.space_before = Pt(8)

    # =============================================================
    # SLIDE 8: TABELA COMPARATIVA
    # =============================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, C_LIGHT_BG)
    add_header(s8, 'Quadro Comparativo: Planilha Matriz vs. Novo Sistema', 'BENEFÍCIOS REAIS', 'Resumo das melhorias operacionais, ganhos de tempo e eliminação de riscos.')

    rows = [
        ('Funcionalidade', 'Planilha Atual (MATRIZ.xlsx)', 'Novo Sistema Web & App'),
        ('Cálculo de Rateio', 'Gera frações decimais; arredondamento manual.', 'Rateio inteligente com compensação automática.'),
        ('Segurança de Fórmulas', 'Alto risco de apagar linhas ou quebrar vínculos.', 'Regras 100% blindadas e seguras no servidor.'),
        ('Cadastro de Produtos', 'Redigitação manual de código e descrição.', 'Catálogo centralizado com autocomplete e histórico.'),
        ('Separação no Estoque', 'Impressão em folhas de papel ou uso de notebook.', 'App Mobile dedicado na mão do separador.'),
        ('Análise de Margem', 'Cálculo isolado que exige conferência manual.', 'Auditoria instantânea de preço e teto de compra.'),
        ('Visual & Experiência', 'Tabelas cinzas padrão Excel.', 'Padrão visual App Érica (Emerald Green / Dark).')
    ]

    table_shape = s8.shapes.add_table(7, 3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.9))
    table = table_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(4.45)
    table.columns[2].width = Inches(4.45)

    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11 if r_idx > 0 else 12)
            p.font.bold = (r_idx == 0 or c_idx == 0)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_DARK_BG
                p.font.color.rgb = C_WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if r_idx % 2 == 1 else C_LIGHT_GREEN
                if c_idx == 2:
                    p.font.color.rgb = C_PRIMARY_DARK
                    p.font.bold = True
                elif c_idx == 1:
                    p.font.color.rgb = RGBColor(185, 28, 28)
                else:
                    p.font.color.rgb = C_TEXT_MAIN

    # =============================================================
    # SLIDE 9: CRONOGRAMA DE IMPLANTAÇÃO (5 SEMANAS)
    # =============================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, C_LIGHT_BG)
    add_header(s9, 'Planejamento e Cronograma de Desenvolvimento', 'ENTREGA EM 5 SEMANAS', 'Metodologia ágil com entregas incrementais e validações práticas semanais.')

    steps = [
        ('Semana 1', 'Arquitetura & Design', 'Estruturação do banco de dados, regras fiscais e aplicação do Design System do App Érica.'),
        ('Semana 2', 'Compras & Pedidos', 'Módulo Web de emissão de pedidos, cadastro de fornecedores e geração de PDFs oficiais.'),
        ('Semana 3', 'Fiscal & Rateio 20 Lojas', 'Implementação do simulador de Limite de Preço e algoritmo de rateio inteligente dos clusters A, B e C.'),
        ('Semana 4', 'App Mobile Multifunção', 'Desenvolvimento do aplicativo Android para compras em viagens e conferência no CD.'),
        ('Semana 5', 'Homologação & Go-Live', 'Testes com pedidos reais, treinamento da equipe e ativação oficial em produção.')
    ]

    for idx, (sem, tit, desc) in enumerate(steps):
        x = Inches(0.8 + idx * 2.38)
        y = Inches(2.1)
        c = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.25), Inches(4.5))
        c.fill.solid()
        c.fill.fore_color.rgb = C_WHITE
        c.line.color.rgb = C_PRIMARY_DARK if idx == 4 else C_CARD_BORDER
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.2)
        
        p_num = tf.paragraphs[0]
        p_num.text = sem.upper()
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = C_PRIMARY_DARK
        
        p_t = tf.add_paragraph()
        p_t.text = tit
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_MAIN
        p_t.space_before = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(8)

    # =============================================================
    # SLIDE 10: PROPOSTA COMERCIAL DA VERSÃO COMPLETA (DARK THEME)
    # =============================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, C_DARK_BG)
    add_header(s10, 'Investimento da Versão Completa (Recomendada)', 'PROPOSTA COMERCIAL ROBUSTA', 'Escopo total com Portal Web, App Mobile Multifunção, Servidor em Nuvem e Suporte.', dark=True)

    props = [
        ('Opção 1: À Vista (com 10% Desconto)', 'R$ 5.800,00', [
            '50% de entrada no início do projeto (R$ 2.900)',
            '50% na entrega e homologação final (R$ 2.900)',
            'Economia direta de R$ 700,00',
            'Prioridade máxima no cronograma de desenvolvimento'
        ], C_PRIMARY),
        ('Opção 2: Parcelado em 3x', 'R$ 6.500,00', [
            'Entrada de R$ 2.200,00 no início do projeto',
            '2ª Parcela de R$ 2.150,00 em 30 dias (entrega módulos Web)',
            '3ª Parcela de R$ 2.150,00 em 60 dias (pós Go-Live e App)',
            'Investimento diluído por marcos de entrega'
        ], C_GOLD)
    ]

    for idx, (pname, pval, pitems, pcolor) in enumerate(props):
        x = Inches(0.8 + idx * 5.95)
        y = Inches(1.9)
        c = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.75), Inches(4.9))
        c.fill.solid()
        c.fill.fore_color.rgb = C_DARK_SURFACE
        c.line.color.rgb = pcolor
        
        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = pname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        pv = tf.add_paragraph()
        pv.text = pval
        pv.font.size = Pt(28)
        pv.font.bold = True
        pv.font.color.rgb = pcolor
        pv.space_before = Pt(4)
        pv.space_after = Pt(12)
        
        for item in pitems:
            pi = tf.add_paragraph()
            pi.text = f'✦  {item}'
            pi.font.size = Pt(12)
            pi.font.color.rgb = RGBColor(203, 213, 225)
            pi.space_before = Pt(6)

    output_path = r'c:\Users\Josué\Documents\App Rafael\Apresentacao_App_Compras_Matriz_Versao_Completa.pptx'
    prs.save(output_path)
    print(f'Sucesso: Apresentação Versão Completa salva em {output_path}')

if __name__ == '__main__':
    create_complete_presentation()
